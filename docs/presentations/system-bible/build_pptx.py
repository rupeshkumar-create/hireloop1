#!/usr/bin/env python3
"""Build a minimal, brand-aligned Hireschema System Bible PPTX."""

from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "Hireschema-System-Bible.pptx"
DOWNLOADS = Path.home() / "Downloads" / "Hireschema-System-Bible.pptx"
LOGO = ROOT / "assets" / "hireschema-mark.png"

# DESIGN.md v2
BG = RGBColor(0x14, 0x14, 0x14)
SURFACE = RGBColor(0x1C, 0x1C, 0x1C)
RAISED = RGBColor(0x24, 0x24, 0x24)
BORDER = RGBColor(0x2A, 0x2A, 0x2A)
TEXT = RGBColor(0xFA, 0xFA, 0xFA)
MUTED = RGBColor(0xA3, 0xA3, 0xA3)
SUBTLE = RGBColor(0x6B, 0x6B, 0x6B)
ACCENT = RGBColor(0xB9, 0xF8, 0x4C)
ACCENT_FG = RGBColor(0x0F, 0x14, 0x00)
DANGER = RGBColor(0xF7, 0x6D, 0x6D)

# Layout grid (16:9)
W = Inches(13.333)
H = Inches(7.5)
ML = Inches(0.85)  # margin left
MR = Inches(0.85)
MT = Inches(0.55)
CONTENT_W = W - ML - MR


def make_logo() -> Path:
    """Rasterize the Hireschema mark (skewed H bars) to PNG."""
    LOGO.parent.mkdir(parents=True, exist_ok=True)
    size = 512
    img = Image.new("RGBA", (size, size), (20, 20, 20, 255))
    draw = ImageDraw.Draw(img)
    lime = (185, 248, 76, 255)
    # Approximate mark bars from hireschema-mark.svg (skewed H)
    # Base rects then apply skew via affine
    mark = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    md = ImageDraw.Draw(mark)
    # scale viewBox 48 → 512
    s = size / 48

    def rect(x, y, w, h):
        md.rectangle(
            [x * s, y * s, (x + w) * s, (y + h) * s],
            fill=lime,
        )

    rect(10.5, 9, 7.5, 12.5)
    rect(10.5, 26.5, 7.5, 12.5)
    rect(30, 9, 7.5, 12.5)
    rect(30, 26.5, 7.5, 12.5)
    rect(10.5, 20.5, 27, 7)
    # skewX(-10°) around center
    skewed = mark.transform(
        (size, size),
        Image.Transform.AFFINE,
        (1, 0.1763, -0.1763 * size / 2, 0, 1, 0),
        resample=Image.Resampling.BICUBIC,
    )
    img.alpha_composite(skewed)
    img.save(LOGO)
    return LOGO


def set_run(run, size: int, *, bold: bool = False, color=TEXT):
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # Force Arial on complex script too
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is None:
            el = rPr.makeelement(qn(f"a:{tag}"), {})
            rPr.append(el)
        el.set("typeface", "Arial")


def solid_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def textbox(slide, left, top, width, height, text, *, size=16, bold=False, color=TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    p.clear()
    run = p.add_run()
    run.text = text
    set_run(run, size, bold=bold, color=color)
    return box


def add_bg(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_logo_lockup(slide, *, mark_size=0.42, word=True):
    slide.shapes.add_picture(str(LOGO), ML, MT, height=Inches(mark_size))
    if word:
        textbox(
            slide,
            ML + Inches(mark_size) + Inches(0.18),
            MT + Inches(0.02),
            Inches(3.5),
            Inches(0.4),
            "Hireschema",
            size=16,
            bold=True,
            color=TEXT,
            anchor=MSO_ANCHOR.MIDDLE,
        )


def page_num(slide, i: int, n: int):
    textbox(
        slide,
        W - MR - Inches(1.1),
        H - Inches(0.45),
        Inches(1.1),
        Inches(0.3),
        f"{i} / {n}",
        size=11,
        color=SUBTLE,
        align=PP_ALIGN.RIGHT,
    )


def title_slide(prs, eyebrow: str, title: str, lede: str, i: int, n: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo_lockup(slide, mark_size=0.48)
    textbox(slide, ML, Inches(2.55), CONTENT_W, Inches(0.35), eyebrow.upper(), size=12, bold=True, color=ACCENT)
    # multiline title support
    box = slide.shapes.add_textbox(ML, Inches(2.95), CONTENT_W, Inches(2.0))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(title.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        set_run(run, 36, bold=True, color=TEXT)
    if lede:
        textbox(slide, ML, Inches(5.2), Inches(10.5), Inches(1.0), lede, size=15, color=MUTED)
    page_num(slide, i, n)
    return slide


def content_slide(prs, eyebrow: str, title: str, bullets: list[str], i: int, n: int, lede: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo_lockup(slide)
    y = Inches(1.15)
    textbox(slide, ML, y, CONTENT_W, Inches(0.3), eyebrow.upper(), size=11, bold=True, color=ACCENT)
    y = Inches(1.45)
    textbox(slide, ML, y, CONTENT_W, Inches(0.65), title, size=26, bold=True, color=TEXT)
    y = Inches(2.2)
    if lede:
        textbox(slide, ML, y, CONTENT_W, Inches(0.55), lede, size=14, color=MUTED)
        y = Inches(2.75)
    # bullets
    box = slide.shapes.add_textbox(ML, y, CONTENT_W, H - y - Inches(0.6))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(bullets[:6]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_before = Pt(2)
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = f"•  {item}"
        set_run(run, 15, color=MUTED)
    page_num(slide, i, n)
    return slide


def two_card_slide(prs, eyebrow: str, title: str, cards: list[tuple[str, str]], i: int, n: int, lede: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo_lockup(slide)
    textbox(slide, ML, Inches(1.15), CONTENT_W, Inches(0.3), eyebrow.upper(), size=11, bold=True, color=ACCENT)
    textbox(slide, ML, Inches(1.45), CONTENT_W, Inches(0.6), title, size=26, bold=True, color=TEXT)
    y = Inches(2.2)
    if lede:
        textbox(slide, ML, y, CONTENT_W, Inches(0.5), lede, size=14, color=MUTED)
        y = Inches(2.75)
    gap = Inches(0.3)
    card_w = (CONTENT_W - gap) / 2
    card_h = Inches(3.2)
    for idx, (label, body) in enumerate(cards[:2]):
        left = ML + idx * (card_w + gap)
        solid_rect(slide, left, y, card_w, card_h, SURFACE, line=BORDER)
        # accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y, Inches(0.08), card_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        textbox(slide, left + Inches(0.35), y + Inches(0.35), card_w - Inches(0.55), Inches(0.35), label.upper(), size=11, bold=True, color=ACCENT)
        textbox(slide, left + Inches(0.35), y + Inches(0.85), card_w - Inches(0.55), card_h - Inches(1.2), body, size=15, color=MUTED)
    page_num(slide, i, n)
    return slide


def four_tile_slide(prs, eyebrow: str, title: str, tiles: list[tuple[str, str, str]], i: int, n: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo_lockup(slide)
    textbox(slide, ML, Inches(1.15), CONTENT_W, Inches(0.3), eyebrow.upper(), size=11, bold=True, color=ACCENT)
    textbox(slide, ML, Inches(1.45), CONTENT_W, Inches(0.55), title, size=26, bold=True, color=TEXT)
    gap = Inches(0.25)
    tile_w = (CONTENT_W - 3 * gap) / 4
    tile_h = Inches(3.6)
    y = Inches(2.3)
    for idx, (path, role, stack) in enumerate(tiles[:4]):
        left = ML + idx * (tile_w + gap)
        solid_rect(slide, left, y, tile_w, tile_h, SURFACE, line=ACCENT)
        textbox(slide, left + Inches(0.25), y + Inches(0.4), tile_w - Inches(0.45), Inches(0.45), path, size=18, bold=True, color=ACCENT)
        textbox(slide, left + Inches(0.25), y + Inches(1.1), tile_w - Inches(0.45), Inches(1.0), role, size=14, bold=True, color=TEXT)
        textbox(slide, left + Inches(0.25), y + Inches(2.2), tile_w - Inches(0.45), Inches(1.0), stack, size=13, color=MUTED)
    page_num(slide, i, n)
    return slide


def build():
    make_logo()
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slides_meta: list = []

    def track(fn, *args, **kwargs):
        slides_meta.append((fn, args, kwargs))

    # Curated minimal content
    track(title_slide, "System bible · Full leadership pack", "Business + candidate\n+ recruiter — one deck", "Product, stack, ingest, match math, agents, and candid status.")
    track(
        two_card_slide,
        "One line",
        "Not a job board with a chatbot",
        [
            ("Aarya · Candidate", "Profile · chat/voice · Google Jobs ingest · scored matches · Request Intro"),
            ("Nitya · Recruiter", "JD · publish · score candidates · in-app intro to candidates · Gmail for HM cold path"),
        ],
        lede="Both share one Supabase graph. Intros handshake via DB — never agent-to-agent RPC.",
    )
    track(
        two_card_slide,
        "Value",
        "Why it exists",
        [
            ("Candidate", "Less apply-spam. Scored roles from live ingest. Warm intro from their Gmail."),
            ("Recruiter", "JD in minutes. Shortlist from the same matcher. Warm intros, not spray."),
        ],
        lede="MVP market: India · INR / LPA · discovery opt-in default OFF",
    )
    track(
        four_tile_slide,
        "Surfaces",
        "What ships where",
        [
            ("web/", "Marketing", "Next.js 15 · Vercel"),
            ("app/", "Candidate + Recruiter SPA", "Zustand · Realtime · SSE"),
            ("api/", "Agents + matching", "FastAPI · LangGraph"),
            ("supabase/", "Shared graph", "Postgres · pgvector · RLS"),
        ],
    )
    track(
        content_slide,
        "Stack",
        "Confirmed technology",
        [
            "Frontend: Next.js 15 · TypeScript · Tailwind · shadcn",
            "Auth: Supabase · LinkedIn OAuth",
            "API: FastAPI · LangGraph · asyncpg",
            "AI: OpenRouter · Claude Sonnet · Gemini Flash · embedding-3-small",
            "Jobs / JD: Apify Google Jobs · Firecrawl",
            "Email: Resend (transactional) · Gmail OAuth (cold intros)",
        ],
    )
    track(
        content_slide,
        "Architecture",
        "How surfaces connect",
        [
            "Apps → FastAPI /api/v1 (auth, matching, ingest, intros)",
            "Aarya: LangGraph tools · writes agent_actions → Realtime UI",
            "Nitya: LISTEN intro_requests · Gmail send after approve-send",
            "Data: Supabase Postgres + pgvector · Auth · Storage · Realtime",
            "Infra path: Railway today · AWS ECS target",
        ],
    )
    track(
        content_slide,
        "Agents",
        "Master loop",
        [
            "Think → choose tool → execute Python → write state + agent_actions",
            "Same pattern for Aarya and Nitya",
            "Intros: Request Intro → intro_requests → NOTIFY → Nitya",
            "Hard rule: no Aarya ↔ Nitya HTTP / RPC",
        ],
    )
    track(title_slide, "Part 2", "Candidate · Aarya\nin detail", "Journey, ingest, match math, Request Intro.")
    track(
        content_slide,
        "Candidate journey",
        "A → Z",
        [
            "LinkedIn OAuth → onboarding (CV · prefs · market)",
            "Chat with Aarya (text or voice on the same surface)",
            "Scored matches feed — India-eligible jobs",
            "Act: Request Intro · Direct Apply · Save",
            "Approve-send cold HM intro from candidate Gmail",
        ],
    )
    track(
        content_slide,
        "Candidate app",
        "SPA + key tools",
        [
            "Routes: /onboarding · /chat · /matches · /profile · /r/{slug}",
            "Client calls FastAPI only (api.ts) — no direct DB writes",
            "Tools: job_search · get_match_score · request_intro · profile_read",
            "Also: save_job · direct_apply · analyze_resume · application kit",
        ],
    )
    track(
        content_slide,
        "Profile graph",
        "Where signals come from",
        [
            "LinkedIn → users / candidates identity fields",
            "Resume → Supabase Storage → parser → skills / experience",
            "Aarya chat writes prefs, location, seniority, must-haves",
            "Embeddings: profile · skills · resume (1536-d via OpenRouter)",
        ],
    )
    track(
        content_slide,
        "Job ingest",
        "Where jobs come from",
        [
            "Apify Google Jobs only — actor johnvc/Google-Jobs-Scraper",
            "Triggers: cron · Aarya auto-ingest · Find new · career path",
            "Pipeline: scrape → upsert jobs → embed → score",
            "Dedup: apify_job_id → fingerprint → apply_url · 24h query skip",
            "Embeddings run in background workers after ingest",
        ],
    )
    track(
        content_slide,
        "Match scoring",
        "overall_score formula",
        [
            "overall = Σ(w × dim) / Σw  then × role · domain · title gates",
            "Weights: Skills 0.40 · Profile 0.30 · Experience 0.15 · Location 0.10 · CTC 0.05",
            "Skills lexical = 0.85×coverage + 0.15×Jaccard (+ embedding lift only)",
            "Missing exp/CTC on the job → those weights are renormalised",
            "Persist floor ~0.35 · feed floor ~0.38",
        ],
        lede="MatchingEngine._assemble_score · services/matching.py",
    )
    track(
        content_slide,
        "Ranking",
        "RRF orders — it doesn’t redefine score",
        [
            "overall_score is computed first and stored on match_scores",
            "RRF (k=60) fuses ranked lists for feed / Aarya search",
            "Feed: fuse overall + skills → MMR · company cap 2",
            "Same MatchingEngine also powers recruiter shortlist",
        ],
    )
    track(
        two_card_slide,
        "Connection rule · R5",
        "Agents never call each other",
        [
            ("Aarya", "Request Intro writes intro_requests. Does not call Nitya."),
            ("Nitya", "LISTEN → enrich → draft → wait for approve-send via Gmail."),
        ],
        lede="Postgres is the only bridge. Cold email = Gmail. Transactional = Resend.",
    )
    track(
        content_slide,
        "Voice · R8",
        "Voice is the chat",
        [
            "Mic lives in the same input bar — no separate voice product",
            "Deepgram Nova-3 STT → same Aarya LangGraph → Aura TTS",
            "TTS timeout guarded (~10s)",
            "Post-MVP: scheduled 15–20 min screening call (later)",
        ],
    )
    track(title_slide, "Part 3", "Recruiter · Nitya\nin detail", "JD import, same MatchingEngine search, pipeline, intros.")
    track(
        content_slide,
        "Recruiter journey",
        "A → Z",
        [
            "Signup as recruiter → create JD (form or URL import)",
            "Publish → public /r/{slug} apply page",
            "Search candidates with the same MatchingEngine",
            "Pipeline kanban · inbox intros",
            "Recruiter → candidate intro is in-app (not Gmail cold path)",
        ],
    )
    track(
        two_card_slide,
        "JD create",
        "Manual + URL import",
        [
            ("Manual", "Structured form: title, location, remote, must/nice skills, comp, JD."),
            ("Import", "Career-page URL via Firecrawl · Ashby/Greenhouse · JSON-LD / HTML."),
        ],
    )
    track(
        content_slide,
        "Graph search",
        "No separate graph ranker",
        [
            "ensure_role_scoring_job → jobs row with source='recruiter'",
            "MatchingEngine.score_job → ORDER BY overall_score",
            "Results land in role_pipeline (stage=search)",
            "Only opted-in candidates: share_with_recruiters · not private",
            "Same weights/gates as the candidate Matches feed",
        ],
    )
    track(
        content_slide,
        "Nitya intros",
        "Three directions — clear roles",
        [
            "Recruiter → candidate: in-app request · candidate accepts in inbox",
            "Candidate → recruiter: in-app on recruiter inbox",
            "Candidate → external HM: Nitya drafts · approve-send · Gmail OAuth",
            "HM path only: Nitya LISTEN + Apify enrichment + Gmail send",
        ],
    )
    track(
        content_slide,
        "Data",
        "Core tables",
        [
            "users · candidates · recruiters · jobs · roles",
            "match_scores · candidate_embeddings · job_embeddings",
            "intro_requests · role_pipeline · agent_actions",
            "gmail_tokens · consent_log · background_jobs",
        ],
    )
    track(
        content_slide,
        "Trust",
        "DPDP + hard rules",
        [
            "consent_log on collection · bias_audit on every match_scores row",
            "GET /me/export · DELETE /me soft-delete · privacy@hireschema.com",
            "Resend = transactional only · Gmail = cold intros after approve-send",
            "No agent RPC · no hardcoded secrets · RLS · sharing default OFF",
        ],
    )
    track(
        two_card_slide,
        "Status & risks",
        "Candid across product",
        [
            ("Candidate", "Match / Find new reliability · chat find-jobs polish · need embed workers after ingest."),
            ("Recruiter", "Depends on role→job sync + embeddings · Gmail OAuth must be live · empty graph → empty shortlist."),
        ],
        lede="Deps: Apify · Gmail OAuth · Resend · durable workers. Payments deferred.",
    )
    track(
        content_slide,
        "Near-term",
        "Next 2 weeks",
        [
            "Harden ingest → embed → score chain",
            "Impression hygiene · Find new reliability",
            "Chat ↔ Matches parity for job search",
            "Gmail intro E2E · embed-after-role-sync for recruiters",
        ],
    )
    track(
        two_card_slide,
        "Post-MVP",
        "Aarya 15–20 min screening call",
        [
            ("For the candidate", "Richer profile → better MatchingEngine input."),
            ("For the recruiter", "On apply: call insights + JD → screening card on pipeline."),
        ],
        lede="After MVP is stable. Builds on today’s Deepgram voice — does not replace score search.",
    )
    track(
        content_slide,
        "Remember",
        "Takeaways",
        [
            "Two agents, one Postgres graph",
            "Jobs from Apify Google Jobs → embed → score",
            "One MatchingEngine for candidate feed and recruiter search",
            "Cold email = Gmail + approve-send · Resend stays transactional",
            "Screening call is post-MVP — after match/ingest stabilize",
        ],
    )

    n = len(slides_meta)
    for i, (fn, args, kwargs) in enumerate(slides_meta, start=1):
        fn(prs, *args, i=i, n=n, **kwargs)

    prs.save(OUT)
    DOWNLOADS.write_bytes(OUT.read_bytes())
    print(f"slides={n}")
    print(f"wrote={OUT}")
    print(f"wrote={DOWNLOADS}")
    print(f"kb={DOWNLOADS.stat().st_size // 1024}")


if __name__ == "__main__":
    build()
